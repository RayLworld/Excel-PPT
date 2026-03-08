import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import threading
import sys
import os
import tempfile
import shutil
import subprocess
import time
import win32com.client
import pandas as pd
from pptx import Presentation

# 自动安装缺失的依赖
def install_package(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package], 
                         creationflags=subprocess.CREATE_NO_WINDOW)

# 检查并安装psutil
try:
    import psutil
except ImportError:
    print("正在安装psutil库...")
    install_package("psutil")
    import psutil

class PPTReplaceGUI:
    def __init__(self, root):
        # 主窗口设置
        self.root = root
        self.root.title("PPT批量替换工具（Excel数据源）")
        self.root.geometry("800x850")
        self.root.resizable(False, False)
        
        # 全局变量
        self.df = None
        self.column_names = []
        self.selected_fields = []
        self.processing = False
        self.app_pid = None  # 存储当前创建的WPS/PPT进程ID
        
        # ========== 1. 路径选择区域 ==========
        frame_path = tk.Frame(root, padx=20, pady=10)
        frame_path.pack(fill=tk.X)
        
        # PPT文件选择
        tk.Label(frame_path, text="PPT模板文件：", font=("微软雅黑", 10)).grid(row=0, column=0, sticky="w", pady=5)
        self.ppt_path_var = tk.StringVar()
        entry_ppt = tk.Entry(frame_path, textvariable=self.ppt_path_var, width=50, font=("微软雅黑", 9))
        entry_ppt.grid(row=0, column=1, padx=10, pady=5)
        btn_ppt = tk.Button(frame_path, text="选择文件", command=self.select_ppt, width=10, font=("微软雅黑", 9))
        btn_ppt.grid(row=0, column=2, padx=5, pady=5)
        
        # Excel文件选择
        tk.Label(frame_path, text="Excel数据文件：", font=("微软雅黑", 10)).grid(row=1, column=0, sticky="w", pady=5)
        self.excel_path_var = tk.StringVar()
        entry_excel = tk.Entry(frame_path, textvariable=self.excel_path_var, width=50, font=("微软雅黑", 9))
        entry_excel.grid(row=1, column=1, padx=10, pady=5)
        btn_excel = tk.Button(frame_path, text="选择文件", command=self.select_excel, width=10, font=("微软雅黑", 9))
        btn_excel.grid(row=1, column=2, padx=5, pady=5)
        
        # 输出目录选择
        tk.Label(frame_path, text="图片输出目录：", font=("微软雅黑", 10)).grid(row=2, column=0, sticky="w", pady=5)
        self.output_dir_var = tk.StringVar()
        entry_output = tk.Entry(frame_path, textvariable=self.output_dir_var, width=50, font=("微软雅黑", 9))
        entry_output.grid(row=2, column=1, padx=10, pady=5)
        btn_output = tk.Button(frame_path, text="选择目录", command=self.select_output_dir, width=10, font=("微软雅黑", 9))
        btn_output.grid(row=2, column=2, padx=5, pady=5)
        
        # ========== 2. Excel参数设置区域 ==========
        frame_excel = tk.Frame(root, padx=20, pady=10)
        frame_excel.pack(fill=tk.X)
        
        # 表头行号设置
        tk.Label(frame_excel, text="Excel表头行号：", font=("微软雅黑", 10)).grid(row=0, column=0, sticky="w", pady=5)
        self.header_row_var = tk.StringVar(value="1")
        entry_header = tk.Entry(frame_excel, textvariable=self.header_row_var, width=10, font=("微软雅黑", 9))
        entry_header.grid(row=0, column=1, padx=10, pady=5)
        tk.Label(frame_excel, text="（按Excel显示的行号，如第2行输入2）", font=("微软雅黑", 9), fg="gray").grid(row=0, column=2, sticky="w", pady=5)
        
        # 读取Excel按钮
        btn_read_excel = tk.Button(frame_excel, text="读取Excel数据", command=self.read_excel, width=15, font=("微软雅黑", 9), bg="#4CAF50", fg="white")
        btn_read_excel.grid(row=0, column=3, padx=20, pady=5)
        
        # ========== 3. 字段选择区域 ==========
        frame_fields = tk.Frame(root, padx=20, pady=10)
        frame_fields.pack(fill=tk.X)
        
        tk.Label(frame_fields, text="选择要替换的字段：", font=("微软雅黑", 10)).grid(row=0, column=0, sticky="w", pady=5)
        
        # 字段列表框
        self.field_listbox = tk.Listbox(frame_fields, selectmode=tk.MULTIPLE, width=70, height=6, font=("微软雅黑", 9))
        self.field_listbox.grid(row=1, column=0, columnspan=2, pady=5)
        
        # 确认选择按钮
        btn_confirm_fields = tk.Button(frame_fields, text="确认选择", command=self.confirm_fields, width=15, font=("微软雅黑", 9), bg="#2196F3", fg="white")
        btn_confirm_fields.grid(row=1, column=2, padx=10, pady=5)
        
        # ========== 4. 进度显示区域 ==========
        frame_progress = tk.Frame(root, padx=20, pady=10)
        frame_progress.pack(fill=tk.X)
        
        tk.Label(frame_progress, text="处理进度：", font=("微软雅黑", 10)).grid(row=0, column=0, sticky="w", pady=5)
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(frame_progress, variable=self.progress_var, maximum=100, length=500)
        self.progress_bar.grid(row=0, column=1, padx=10, pady=5)
        self.progress_label = tk.Label(frame_progress, text="0%", font=("微软雅黑", 9))
        self.progress_label.grid(row=0, column=2, padx=5, pady=5)
        
        # ========== 5. 日志输出区域 ==========
        frame_log = tk.Frame(root, padx=20, pady=10)
        frame_log.pack(fill=tk.BOTH, expand=True)
        
        tk.Label(frame_log, text="运行日志：", font=("微软雅黑", 10)).grid(row=0, column=0, sticky="w", pady=5)
        self.log_text = tk.Text(frame_log, width=90, height=15, font=("微软雅黑", 9), state=tk.DISABLED)
        scrollbar = tk.Scrollbar(frame_log, orient=tk.VERTICAL, command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)
        self.log_text.grid(row=1, column=0, sticky="nsew", pady=5)
        scrollbar.grid(row=1, column=1, sticky="ns", pady=5)
        
        # ========== 6. 开始处理按钮 ==========
        frame_btn = tk.Frame(root, padx=20, pady=10)
        frame_btn.pack(fill=tk.X)
        
        self.btn_start = tk.Button(frame_btn, text="开始批量替换", command=self.start_process, width=20, font=("微软雅黑", 10), bg="#FF9800", fg="white")
        self.btn_start.grid(row=0, column=0, pady=10)
        
        # 初始化日志
        self.log("✅ 程序已启动，请按步骤操作：")
        self.log("1. 选择PPT模板、Excel数据文件、输出目录")
        self.log("2. 设置Excel表头行号，点击【读取Excel数据】")
        self.log("3. 在字段列表中选择要替换的字段，点击【确认选择】")
        self.log("4. 点击【开始批量替换】执行操作\n")

    # ========== 核心功能函数 ==========
    def log(self, msg):
        """日志输出到文本框"""
        self.log_text.config(state=tk.NORMAL)
        self.log_text.insert(tk.END, f"{time.strftime('%H:%M:%S')} - {msg}\n")
        self.log_text.see(tk.END)
        self.log_text.config(state=tk.DISABLED)
        self.root.update_idletasks()

    def select_ppt(self):
        """选择PPT文件"""
        path = filedialog.askopenfilename(
            title="选择PPT模板文件",
            filetypes=[("PPT文件", "*.pptx;*.ppt"), ("所有文件", "*.*")]
        )
        if path:
            self.ppt_path_var.set(os.path.abspath(path))
            self.log(f"📄 已选择PPT文件：{path}")

    def select_excel(self):
        """选择Excel文件"""
        path = filedialog.askopenfilename(
            title="选择Excel数据文件",
            filetypes=[("Excel文件", "*.xlsx;*.xls"), ("所有文件", "*.*")]
        )
        if path:
            self.excel_path_var.set(os.path.abspath(path))
            self.log(f"📄 已选择Excel文件：{path}")

    def select_output_dir(self):
        """选择输出目录"""
        path = filedialog.askdirectory(title="选择图片输出目录")
        if path:
            self.output_dir_var.set(os.path.abspath(path))
            self.log(f"📂 已选择输出目录：{path}")

    def read_excel(self):
        """读取Excel数据"""
        # 验证输入
        excel_path = self.excel_path_var.get().strip()
        if not excel_path:
            messagebox.showerror("错误", "请先选择Excel文件！")
            return
        
        try:
            header_row = int(self.header_row_var.get().strip())
            if header_row < 1:
                messagebox.showerror("错误", "表头行号必须≥1！")
                return
        except ValueError:
            messagebox.showerror("错误", "表头行号必须是数字！")
            return
        
        # 读取Excel
        self.log(f"📖 开始读取Excel数据（表头行：{header_row}）...")
        try:
            pandas_header = header_row - 1
            df = pd.read_excel(excel_path, header=pandas_header)
            df = df.dropna(how='all').reset_index(drop=True)
            df = df.dropna(axis=1, how='all')
            column_names = list(df.columns)
            column_names = [col if not pd.isna(col) else f"列_{i+1}" for i, col in enumerate(column_names)]
            
            self.df = df
            self.column_names = column_names
            
            # 清空并填充字段列表框
            self.field_listbox.delete(0, tk.END)
            for idx, col in enumerate(column_names):
                self.field_listbox.insert(idx, col)
            
            self.log(f"✅ Excel读取成功！共{len(df)}行有效数据，包含字段：{', '.join(column_names)}")
            messagebox.showinfo("成功", f"Excel读取成功！\n共{len(df)}行数据，{len(column_names)}个字段")
        except Exception as e:
            self.log(f"❌ Excel读取失败：{str(e)}")
            messagebox.showerror("错误", f"Excel读取失败：{str(e)}")

    def confirm_fields(self):
        """确认选择的替换字段"""
        if not self.column_names:
            messagebox.showerror("错误", "请先读取Excel数据！")
            return
        
        # 获取选中的字段
        selected_indices = self.field_listbox.curselection()
        if not selected_indices:
            messagebox.showwarning("警告", "请至少选择一个替换字段！")
            return
        
        self.selected_fields = [self.column_names[idx] for idx in selected_indices]
        self.log(f"🎯 已确认替换字段：{', '.join(self.selected_fields)}")
        messagebox.showinfo("成功", f"已选择以下字段：\n{', '.join(self.selected_fields)}")

    def format_number(self, value):
        """格式化数字，去除无效小数点"""
        if isinstance(value, (int, float)):
            if value.is_integer():
                return str(int(value))
            else:
                return str(value)
        else:
            return str(value)
    
    def get_process_pid(self, app):
        """获取COM对象对应的进程ID"""
        try:
            # 不同COM对象获取PID的方式
            if hasattr(app, 'HWND'):
                hwnd = app.HWND
                # 通过窗口句柄获取进程ID
                import win32process
                pid = win32process.GetWindowThreadProcessId(hwnd)[1]
                return pid
            elif hasattr(app, 'Application'):
                return self.get_process_pid(app.Application)
            else:
                return None
        except:
            return None
    
    def kill_specific_process(self, pid, process_names):
        """精准关闭指定PID的进程"""
        if not pid:
            return
        
        try:
            # 先检查进程是否存在
            if psutil.pid_exists(pid):
                process = psutil.Process(pid)
                # 验证进程名称是否匹配（防止误杀）
                if process.name().lower() in [name.lower() for name in process_names]:
                    # 先正常终止，失败再强制终止
                    process.terminate()
                    time.sleep(1)
                    if process.is_running():
                        process.kill()
                    self.log(f"✅ 已关闭进程 {pid} ({process.name()})")
                else:
                    self.log(f"⚠️  进程PID {pid} 名称不匹配，跳过关闭")
            else:
                self.log(f"⚠️  进程PID {pid} 已不存在")
        except Exception as e:
            self.log(f"⚠️  关闭进程 {pid} 失败：{str(e)}")

    def replace_ppt_by_row(self, row_idx):
        """按行替换PPT并导出图片（后台静默运行，无界面弹出）"""
        ppt_path = self.ppt_path_var.get()
        output_dir = self.output_dir_var.get()
        
        # 重置进程ID
        self.app_pid = None
        
        # 创建临时PPT
        temp_dir = tempfile.gettempdir()
        temp_ppt_name = f"temp_{os.urandom(4).hex()}.pptx"
        temp_ppt_path = os.path.join(temp_dir, temp_ppt_name)
        shutil.copy2(ppt_path, temp_ppt_path)
        
        # 替换文本
        prs = Presentation(temp_ppt_path)
        row_data = self.df.iloc[row_idx]
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for field in self.selected_fields:
                                if field in run.text:
                                    raw_value = row_data[field]
                                    replace_value = self.format_number(raw_value)
                                    run.text = run.text.replace(field, replace_value)
        
        prs.save(temp_ppt_path)
        prs = None
        
        # ========== 核心逻辑：后台静默运行，无界面弹出 ==========
        app = None
        presentation = None
        slide_count = 0
        is_use_wps = True
        
        try:
            # 第一步：尝试所有WPS演示类名（静默模式）
            wpp_class_list = [
                "KWPP.Application.9",
                "KWPP.Application.12",
                "KWPP.Application",
                "WPP.Application",
                "KSO.Application",
            ]
            
            for cls_name in wpp_class_list:
                try:
                    # 创建WPS对象并设置静默模式
                    app = win32com.client.Dispatch(cls_name)
                    if app is not None:
                        # 关键设置1：隐藏主界面
                        app.Visible = False
                        # 关键设置2：禁用显示警报
                        app.DisplayAlerts = False
                        # 关键设置3：WPS特有静默参数（部分版本需要）
                        if hasattr(app, 'Silent'):
                            app.Silent = True
                        if hasattr(app, 'ScreenUpdating'):
                            app.ScreenUpdating = False
                        
                        # 获取并记录进程ID
                        self.app_pid = self.get_process_pid(app)
                        break
                except:
                    continue
            
            # 第二步：如果WPS失败，尝试PowerPoint（静默模式）
            if app is None:
                is_use_wps = False
                try:
                    app = win32com.client.Dispatch("PowerPoint.Application")
                    # PowerPoint静默设置
                    app.Visible = False
                    app.DisplayAlerts = False
                    
                    self.app_pid = self.get_process_pid(app)
                    self.log(f"✅ WPS调用失败，已降级使用Microsoft PowerPoint（静默模式），进程ID：{self.app_pid}")
                except Exception as ppt_e:
                    raise Exception(f"❌ WPS和Microsoft PowerPoint均调用失败！错误：{ppt_e}")
            
            # 导出图片核心逻辑
            presentation = app.Presentations.Open(
                temp_ppt_path,
                ReadOnly=True,       # 只读模式打开，避免锁定文件
                WithWindow=False     # PowerPoint关键：不显示窗口
            )
            slide_count = presentation.Slides.Count
            
            # 获取PPT页面实际尺寸并动态导出
            first_slide = presentation.Slides(1)
            slide_width_pts = first_slide.Parent.PageSetup.SlideWidth
            slide_height_pts = first_slide.Parent.PageSetup.SlideHeight
            slide_width_px = int(slide_width_pts * 96 / 72)
            slide_height_px = int(slide_height_pts * 96 / 72)            
            # 按实际尺寸导出每张幻灯片
            for i in range(slide_count):
                slide = presentation.Slides(i + 1)
                img_name = f"row{row_idx+1}_slide_{i+1}.png"
                img_path = os.path.join(output_dir, img_name)
                slide.Export(img_path, "PNG", slide_width_px, slide_height_px)
                time.sleep(0.2)  # 缩短等待时间，提升效率
        
        except Exception as e:
            self.log(f"❌ 图片导出失败：{str(e)}")
            raise Exception(f"导出图片失败，请确保已安装WPS或Microsoft Office！错误信息：{str(e)}")
        
        finally:
            # 安全清理资源
            if presentation:
                presentation.Close()
            if app:
                app.Quit()
            
            # 精准关闭当前程序创建的进程（不影响其他WPS/PPT）
            if is_use_wps:
                self.kill_specific_process(self.app_pid, ["wps.exe", "wpp.exe", "ksolaunch.exe"])
            else:
                self.kill_specific_process(self.app_pid, ["POWERPNT.EXE"])
            
            # 删除临时文件
            if os.path.exists(temp_ppt_path):
                os.remove(temp_ppt_path)
        
        return slide_count

    def start_process(self):
        """开始批量处理（子线程执行）"""
        # 验证前置条件
        if self.processing:
            messagebox.showwarning("警告", "正在处理中，请等待！")
            return
        
        if not self.ppt_path_var.get():
            messagebox.showerror("错误", "请选择PPT文件！")
            return
        
        if not self.output_dir_var.get():
            messagebox.showerror("错误", "请选择输出目录！")
            return
        
        if self.df is None:
            messagebox.showerror("错误", "请先读取Excel数据！")
            return
        
        if not self.selected_fields:
            messagebox.showerror("错误", "请先选择替换字段！")
            return
        
        # 创建输出目录
        os.makedirs(self.output_dir_var.get(), exist_ok=True)
        
        # 启动子线程
        self.processing = True
        self.btn_start.config(state=tk.DISABLED)
        self.progress_var.set(0)
        self.progress_label.config(text="0%")
        
        threading.Thread(target=self.process_thread, daemon=True).start()

    def process_thread(self):
        """处理线程"""
        total_rows = len(self.df)
        self.log(f"🚀 开始批量处理，共{total_rows}行数据...")
        
        try:
            for row_idx in range(total_rows):
                self.log(f"🔄 正在处理第{row_idx+1}/{total_rows}行数据...")
                slide_count = self.replace_ppt_by_row(row_idx)
                self.log(f"✅ 第{row_idx+1}行处理完成，导出{slide_count}张图片")
                
                # 更新进度
                progress = (row_idx + 1) / total_rows * 100
                self.progress_var.set(progress)
                self.progress_label.config(text=f"{progress:.1f}%")
            
            self.log("\n🎉 所有数据处理完成！")
            self.log(f"📂 图片已保存至：{self.output_dir_var.get()}")
            messagebox.showinfo("完成", "所有数据处理完成！\n已打开输出目录")
            os.startfile(self.output_dir_var.get())
        
        except Exception as e:
            self.log(f"❌ 处理失败：{str(e)}")
            messagebox.showerror("错误", f"处理失败：{str(e)}")
        
        finally:
            self.processing = False
            self.btn_start.config(state=tk.NORMAL)
            self.progress_var.set(100)
            self.progress_label.config(text="100%")

if __name__ == "__main__":
    # 启动GUI
    root = tk.Tk()
    app = PPTReplaceGUI(root)
    root.mainloop()