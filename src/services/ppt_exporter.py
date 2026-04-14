import os
import shutil
import subprocess
import tempfile
import time

import win32com.client
from pptx import Presentation

WPS_CLASS_LIST = [
    "KWPP.Application.9",
    "KWPP.Application.12",
    "KWPP.Application",
    "WPP.Application",
    "KSO.Application",
]


def format_cell_value(value):
    if value is None:
        return ""
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return str(value)
    return str(value)


def _open_presentation_app(log_func):
    app = None
    is_use_wps = True

    for cls_name in WPS_CLASS_LIST:
        try:
            app = win32com.client.Dispatch(cls_name)
            if app is not None:
                break
        except Exception:
            continue

    if app is None:
        is_use_wps = False
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            log_func("✅ WPS调用失败，已降级使用Microsoft PowerPoint")
        except Exception as ppt_error:
            raise Exception(f"❌ WPS和Microsoft PowerPoint均调用失败！错误：{ppt_error}")

    app.DisplayAlerts = False
    return app, is_use_wps


def _cleanup_office_process(is_use_wps):
    if is_use_wps:
        subprocess.run(
            ["taskkill", "/f", "/im", "wps.exe", "/im", "wpp.exe", "/im", "ksolaunch.exe"],
            capture_output=True,
            timeout=5,
            creationflags=subprocess.CREATE_NO_WINDOW,
        )
    else:
        subprocess.run(
            ["taskkill", "/f", "/im", "POWERPNT.EXE"],
            capture_output=True,
            timeout=5,
            creationflags=subprocess.CREATE_NO_WINDOW,
        )


def replace_and_export_row(ppt_path, output_dir, row_data, selected_fields, row_idx, log_func):
    temp_dir = tempfile.gettempdir()
    temp_ppt_name = f"temp_{os.urandom(4).hex()}.pptx"
    temp_ppt_path = os.path.join(temp_dir, temp_ppt_name)
    shutil.copy2(ppt_path, temp_ppt_path)

    presentation_obj = Presentation(temp_ppt_path)
    for slide in presentation_obj.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for field in selected_fields:
                            if field in run.text:
                                replace_value = format_cell_value(row_data.get(field, ""))
                                run.text = run.text.replace(field, replace_value)

    presentation_obj.save(temp_ppt_path)
    presentation_obj = None

    app = None
    opened_presentation = None
    slide_count = 0
    is_use_wps = True

    try:
        app, is_use_wps = _open_presentation_app(log_func)
        opened_presentation = app.Presentations.Open(temp_ppt_path)
        slide_count = opened_presentation.Slides.Count

        for i in range(slide_count):
            slide = opened_presentation.Slides(i + 1)
            img_name = f"row{row_idx + 1}_slide_{i + 1}.png"
            img_path = os.path.join(output_dir, img_name)
            slide.Export(img_path, "PNG", 1920, 1080)
            time.sleep(0.5)
    finally:
        if opened_presentation:
            opened_presentation.Close()
        if app:
            app.Quit()
        _cleanup_office_process(is_use_wps)
        if os.path.exists(temp_ppt_path):
            os.remove(temp_ppt_path)

    return slide_count
